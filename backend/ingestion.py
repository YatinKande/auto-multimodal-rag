
import os
import io
import pandas as pd
from typing import List, Dict, Any
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import ollama
from backend.config import Config

class IngestionEngine:
    def __init__(self):
        self.chunk_size = 500
        self.chunk_overlap = 50

    def process_file(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(filename)[1].lower()
        chunks = []
        
        if ext == ".pdf":
            chunks = self._process_pdf(file_path, filename)
        elif ext == ".docx":
            chunks = self._process_docx(file_path, filename)
        elif ext == ".pptx":
            chunks = self._process_pptx(file_path, filename)
        elif ext == ".txt":
            chunks = self._process_txt(file_path, filename)
        elif ext == ".csv":
            chunks = self._process_csv(file_path, filename)
        elif ext in [".jpg", ".jpeg", ".png"]:
            chunks = self._process_image(file_path, filename)
        
        return chunks

    def _chunk_text(self, text: str, source_info: Dict[str, Any]) -> List[Dict[str, Any]]:
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size
            chunk_text = text[start:end]
            chunks.append({
                "text": chunk_text,
                "metadata": source_info
            })
            start += self.chunk_size - self.chunk_overlap
        return chunks

    def _describe_image(self, image_bytes: bytes) -> str:
        try:
            # Convert bytes to PIL Image to verify it works, then back to bytes if needed
            # Ollama python client accepts bytes directly usually, or path.
            # Let's save to a temp buffer or pass bytes.
            response = ollama.chat(model='llama3.2-vision', messages=[
                {
                    'role': 'user',
                    'content': 'You are an automotive expert. Identify the devices in this diagnostic setup. 1. The handheld device with a screen is likely an "OBD-II Scanner" or "Scan Tool" (NOT a smartphone). 2. Transcribe any text on the scanner screen (e.g. "Mode $03", "OTTC"). 3. Describe the laptop and vehicle connection. Be precise.',
                    'images': [image_bytes]
                }
            ])
            return response['message']['content']
        except Exception as e:
            print(f"Error describing image: {e}")
            return "Image description unavailable."

    def _process_pdf(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        chunks = []
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                chunks.extend(self._chunk_text(text, {"filename": filename, "page": i + 1, "type": "text"}))
            
            # Extract images
            for image_file_object in page.images:
                try:
                    desc = self._describe_image(image_file_object.data)
                    chunks.append({
                        "text": f"[IMAGE DESCRIPTION] {desc}",
                        "metadata": {"filename": filename, "page": i + 1, "type": "image_description"}
                    })
                except Exception as e:
                    print(f"Failed to process image in PDF: {e}")
        return chunks

    def _process_docx(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        chunks = []
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        
        text = "\n".join(full_text)
        chunks.extend(self._chunk_text(text, {"filename": filename, "page": None, "type": "text"}))
        
        # Extract images (basic approach for docx)
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    desc = self._describe_image(image_data)
                    chunks.append({
                        "text": f"[IMAGE DESCRIPTION] {desc}",
                        "metadata": {"filename": filename, "page": None, "type": "image_description"}
                    })
                except:
                    pass
        return chunks

    def _process_pptx(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        chunks = []
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                
                if hasattr(shape, "image"):
                    try:
                        image_data = shape.image.blob
                        desc = self._describe_image(image_data)
                        chunks.append({
                            "text": f"[IMAGE DESCRIPTION] {desc}",
                            "metadata": {"filename": filename, "page": i + 1, "type": "image_description"}
                        })
                    except:
                        pass
            
            if text:
                chunks.extend(self._chunk_text(text, {"filename": filename, "page": i + 1, "type": "text"}))
        return chunks

    def _process_txt(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        with open(file_path, 'r') as f:
            text = f.read()
        return self._chunk_text(text, {"filename": filename, "page": None, "type": "text"})

    def _process_csv(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        df = pd.read_csv(file_path)
        chunks = []
        for idx, row in df.iterrows():
            text = " ".join([f"{col}: {val}" for col, val in row.items()])
            chunks.append({
                "text": text,
                "metadata": {"filename": filename, "row": idx, "type": "csv_row"}
            })
        return chunks

    def _process_image(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        try:
            with open(file_path, "rb") as f:
                image_bytes = f.read()
            desc = self._describe_image(image_bytes)
            return [{
                "text": f"[IMAGE DESCRIPTION] {desc}",
                "metadata": {"filename": filename, "type": "image_description"}
            }]
        except Exception as e:
            print(f"Error processing image file {filename}: {e}")
            return []
